#!/usr/bin/env python3
"""
anki_to_slides.py
Convert a tab-separated Anki plain-text export into a 16:9 slide deck.

Input:
    A .txt file where column 1 is the card front and column 2 is the card back,
    separated by tabs. Basic HTML tags are stripped. ``<img src="...">`` tags are
    resolved against the --media directory. Anki's ``<hr id="answer">`` marker
    is respected: anything preceding it in the back field is discarded so the
    front's content never leaks onto the back slide.

Output (one page/slide per card side, order front1, back1, front2, back2, ...):
    --format pdf   -> export/<stem>.pdf
    --format pptx  -> export/<stem>/<stem>.pptx       (wrapped in a folder)
    --format png   -> export/<stem>/slide_NNN.png     (folder of PNGs)

The per-format output is exposed both via CLI and as importable functions
(``render_pdf``, ``render_pptx``, ``render_png``) so this module can be reused
from a web backend later without changes.

Usage:
    python anki_to_slides.py                              # single .txt in ./import, PDF
    python anki_to_slides.py --format pptx
    python anki_to_slides.py --format png
    python anki_to_slides.py import/deck.txt --format pdf --out export/deck.pdf
"""
from __future__ import annotations

import argparse
import csv
import html
import io
import re
import sys
import zipfile
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path
from typing import List, Optional, Tuple
from urllib.parse import unquote

from PIL import Image, ImageDraw, ImageFont


# ---------------------------------------------------------------------------
# Geometry (all dimensions in inches unless otherwise noted)
# ---------------------------------------------------------------------------
SLIDE_WIDTH_IN = 13.3333
SLIDE_HEIGHT_IN = 7.5
SAFE_MARGIN_IN = 1.0

# When a card has both text and images, the image occupies the top 2/3 of the
# safe area and the text sits in the bottom 1/3.
IMAGE_REGION_FRAC = 2.0 / 3.0
TEXT_REGION_FRAC = 1.0 / 3.0

IMAGE_STACK_GAP_IN = 0.1  # vertical gap between stacked images

MAX_FONT_PT = 54
MIN_FONT_PT = 10
LINE_HEIGHT_FACTOR = 1.25

# PNG rasterization resolution. A 13.33"×7.5" widescreen slide at 240 DPI
# renders to 3200×1800 px — sharp on Retina/HiDPI displays where a 1920-wide
# image otherwise shows at half size and reads as microscopic. Zip size
# roughly doubles vs. the old 144 DPI; acceptable for a format whose whole
# point is "each card as its own shareable image."
PNG_DPI = 240

DEFAULT_IMPORT_DIR = Path("import")
DEFAULT_EXPORT_DIR = Path("export")

# Default Anki media location on macOS. Override with --media if needed.
DEFAULT_ANKI_MEDIA = (
    Path.home()
    / "Library"
    / "Application Support"
    / "Anki2"
    / "User 1"
    / "collection.media"
)

SUPPORTED_FORMATS = ("pdf", "pptx", "png", "apkg", "anki-txt")


# ---------------------------------------------------------------------------
# HTML parsing helpers
# ---------------------------------------------------------------------------
IMG_TAG_RE = re.compile(
    r'<img\b[^>]*?\bsrc\s*=\s*["\']([^"\']+)["\'][^>]*>', re.IGNORECASE
)
TAG_RE = re.compile(r'<[^>]+>')
BR_RE = re.compile(r'<\s*br\s*/?\s*>', re.IGNORECASE)
WS_RUN_RE = re.compile(r'[ \t]+')
BLANK_LINE_RUN_RE = re.compile(r'\n{3,}')

# Anki's back template typically prepends {{FrontSide}}<hr id="answer">{{Back}}.
# Everything up through that marker belongs to the front and must not appear on
# the back slide.
ANKI_ANSWER_SEP_RE = re.compile(
    r'.*<hr\b[^>]*\bid\s*=\s*["\']?answer["\']?[^>]*>',
    re.IGNORECASE | re.DOTALL,
)

# Anki audio/video references like [sound:foo.mp3] — slides can't play these,
# so strip them entirely rather than leaving literal brackets in the text.
SOUND_TAG_RE = re.compile(r'\[sound:[^\]]*\]', re.IGNORECASE)

# Cloze deletions: {{c1::answer}} or {{c1::answer::hint}}. In .txt exports with
# "Include HTML" these are rendered by Anki, but when we reconstruct notes from
# a raw .apkg/.colpkg SQLite we see the source. Unwrap to just the answer so
# the slide shows something sensible instead of {{c1::…}}.
CLOZE_RE = re.compile(r'\{\{c\d+::(.*?)(?:::[^}]*)?\}\}', re.DOTALL)


@dataclass
class CardSide:
    """Parsed representation of one side of an Anki card."""
    text: str
    images: List[Path]


def parse_side(raw: str, media_dir: Path, is_back: bool = False) -> CardSide:
    """Split a raw HTML-ish Anki field into plain text + referenced images."""
    if raw is None:
        return CardSide(text="", images=[])

    if is_back:
        raw = ANKI_ANSWER_SEP_RE.sub("", raw, count=1)

    raw = SOUND_TAG_RE.sub("", raw)
    raw = CLOZE_RE.sub(r"\1", raw)

    images: List[Path] = []

    def _capture_img(match: re.Match) -> str:
        src = match.group(1).strip()
        images.append(media_dir / src)
        return " "

    text = IMG_TAG_RE.sub(_capture_img, raw)
    text = BR_RE.sub("\n", text)
    text = TAG_RE.sub(" ", text)
    text = html.unescape(text)

    lines = [WS_RUN_RE.sub(" ", ln).strip() for ln in text.splitlines()]
    text = "\n".join(lines).strip()
    text = BLANK_LINE_RUN_RE.sub("\n\n", text)
    return CardSide(text=text, images=images)


def read_cards_from_text(text: str, media_dir: Path) -> List[CardSide]:
    """Parse a tab-separated Anki export string into card sides."""
    sides: List[CardSide] = []
    reader = csv.reader(io.StringIO(text), delimiter="\t", quotechar='"')
    for row in reader:
        if not row:
            continue
        if row[0].startswith("#"):
            continue
        front = row[0] if len(row) >= 1 else ""
        back = row[1] if len(row) >= 2 else ""
        sides.append(parse_side(front, media_dir, is_back=False))
        sides.append(parse_side(back, media_dir, is_back=True))
    return sides


def read_cards(path: Path, media_dir: Path) -> List[CardSide]:
    """Read an Anki TSV export file and return sides in order front1, back1, ..."""
    with open(path, "r", encoding="utf-8", newline="") as fh:
        return read_cards_from_text(fh.read(), media_dir)


# ---------------------------------------------------------------------------
# Notion parsing
#
# Notion exports collapsible "toggle" blocks as ``<details><summary>...</summary>
# ...children...</details>`` in both its HTML and its Markdown export. The
# Markdown export embeds the same HTML inline, so the two share one parser:
# find every ``<details>`` element, take its first ``<summary>`` child as the
# card front, take everything else inside as the card back. Per the user's
# spec we treat each toggle (at any depth) as its own card; the parent's
# back HTML has any nested ``<details>...</details>`` removed before
# rendering so its inner toggles don't get duplicated as back text.
#
# We only emit cards from toggle blocks. Non-toggle Notion content (plain
# headings, paragraphs, tables, callouts, etc.) is intentionally ignored —
# users who want to see that content should keep using Notion.
# ---------------------------------------------------------------------------


class _ToggleCollector(HTMLParser):
    """Walk an HTML document and collect every ``<details>`` element's raw HTML.

    For each top-level-or-nested ``<details>`` we record:
      - ``summary_html``: the inner HTML of its first ``<summary>`` child
      - ``body_html``: the rest of the children's inner HTML, with any nested
        ``<details>...</details>`` stripped (those become their own cards via
        a separate pass over the same text).

    We work with raw HTML strings (rather than building a DOM) so the
    downstream ``parse_side`` cleanup — which already understands ``<img>``,
    ``<br>``, cloze, sound tags — applies unchanged.
    """

    def __init__(self) -> None:
        super().__init__(convert_charrefs=False)
        self._details_stack: List[dict] = []
        self.toggles: List[Tuple[str, str]] = []  # (summary_html, body_html)

    def _emit_raw(self, raw: str) -> None:
        if not self._details_stack:
            return
        frame = self._details_stack[-1]
        if frame["in_summary"] and frame["summary_depth"] == 0:
            frame["summary"].append(raw)
        else:
            frame["body"].append(raw)

    def handle_starttag(self, tag: str, attrs) -> None:
        raw = self.get_starttag_text() or f"<{tag}>"
        tag_l = tag.lower()
        if tag_l == "details":
            self._details_stack.append({
                "summary": [],
                "body": [],
                "in_summary": False,
                "summary_depth": 0,
            })
            return
        if not self._details_stack:
            return
        frame = self._details_stack[-1]
        if tag_l == "summary" and not frame["in_summary"]:
            frame["in_summary"] = True
            frame["summary_depth"] = 0
            return
        if frame["in_summary"]:
            if tag_l == "summary":
                frame["summary_depth"] += 1
            frame["summary"].append(raw)
        else:
            frame["body"].append(raw)

    def handle_startendtag(self, tag: str, attrs) -> None:
        raw = self.get_starttag_text() or f"<{tag}/>"
        if tag.lower() == "details":
            return
        self._emit_raw(raw)

    def handle_endtag(self, tag: str) -> None:
        tag_l = tag.lower()
        if tag_l == "details":
            if not self._details_stack:
                return
            frame = self._details_stack.pop()
            summary_html = "".join(frame["summary"]).strip()
            body_html = "".join(frame["body"]).strip()
            self.toggles.append((summary_html, body_html))
            # Re-emit the closed nested details into the parent's body as a
            # placeholder marker we can later strip — keeps the parent's
            # back from accidentally containing the child's contents.
            if self._details_stack:
                self._details_stack[-1]["body"].append("")
            return
        if not self._details_stack:
            return
        frame = self._details_stack[-1]
        if frame["in_summary"]:
            if tag_l == "summary":
                if frame["summary_depth"] == 0:
                    frame["in_summary"] = False
                    return
                frame["summary_depth"] -= 1
            frame["summary"].append(f"</{tag}>")
        else:
            frame["body"].append(f"</{tag}>")

    def handle_data(self, data: str) -> None:
        self._emit_raw(data)

    def handle_entityref(self, name: str) -> None:
        self._emit_raw(f"&{name};")

    def handle_charref(self, name: str) -> None:
        self._emit_raw(f"&#{name};")

    def handle_comment(self, data: str) -> None:
        self._emit_raw(f"<!--{data}-->")


# Notion's markdown export embeds raw HTML for toggles, so we don't need a
# full markdown parser — we only need to (a) turn `![alt](path)` image
# references into `<img src="path">` so the downstream parser sees them,
# and (b) leave the rest of the source alone (the HTML walker will simply
# treat it as text data, which gets stripped by parse_side).
_MD_IMAGE_RE = re.compile(r'!\[([^\]]*)\]\(([^)\s]+)(?:\s+"[^"]*")?\)')


def _markdown_images_to_html(md: str) -> str:
    return _MD_IMAGE_RE.sub(
        lambda m: f'<img src="{html.escape(m.group(2), quote=True)}" alt="{html.escape(m.group(1), quote=True)}">',
        md,
    )


def _resolve_image_path(src: str, media_dir: Path) -> Path:
    """Resolve an ``<img src>`` reference to a local file under ``media_dir``.

    Notion percent-encodes spaces and unicode in image paths and stores
    images in subfolders like ``Page Title abc123/``. We try the path as
    written, the percent-decoded path, and the basename of either, all
    relative to ``media_dir``. The first one that exists wins; otherwise
    we return ``media_dir / basename`` so the renderer's "image not found"
    warning surfaces a useful filename.
    """
    candidates: List[Path] = []
    for s in (src, unquote(src)):
        candidates.append(media_dir / s)
        candidates.append(media_dir / Path(s).name)
    for c in candidates:
        try:
            if c.is_file():
                return c
        except OSError:
            continue
    return media_dir / Path(unquote(src)).name


def _parse_side_with_resolver(raw: str, media_dir: Path, is_back: bool = False) -> CardSide:
    """Like ``parse_side`` but resolves image paths against Notion's nested folders.

    The standard ``parse_side`` joins ``media_dir / src`` directly; for Notion
    exports the ``src`` may be percent-encoded or live in a sibling folder.
    We capture each ``<img>`` ourselves with the resolver, then hand the
    text fragments to the same cleanup pipeline.
    """
    if raw is None:
        return CardSide(text="", images=[])

    if is_back:
        raw = ANKI_ANSWER_SEP_RE.sub("", raw, count=1)

    raw = SOUND_TAG_RE.sub("", raw)
    raw = CLOZE_RE.sub(r"\1", raw)

    images: List[Path] = []

    def _capture_img(match: re.Match) -> str:
        src = match.group(1).strip()
        images.append(_resolve_image_path(src, media_dir))
        return " "

    text = IMG_TAG_RE.sub(_capture_img, raw)
    text = BR_RE.sub("\n", text)
    text = TAG_RE.sub(" ", text)
    text = html.unescape(text)

    lines = [WS_RUN_RE.sub(" ", ln).strip() for ln in text.splitlines()]
    text = "\n".join(lines).strip()
    text = BLANK_LINE_RUN_RE.sub("\n\n", text)
    return CardSide(text=text, images=images)


def read_cards_from_notion_html(html_text: str, media_dir: Path) -> List[CardSide]:
    """Parse a Notion HTML export string into card sides (front, back, ...)."""
    if not html_text:
        return []
    collector = _ToggleCollector()
    try:
        collector.feed(html_text)
        collector.close()
    except Exception:
        # HTMLParser is forgiving but Notion sometimes ships malformed
        # snippets; surface what we got rather than failing the whole upload.
        pass
    sides: List[CardSide] = []
    for summary_html, body_html in collector.toggles:
        front = _parse_side_with_resolver(summary_html, media_dir, is_back=False)
        back = _parse_side_with_resolver(body_html, media_dir, is_back=False)
        if not front.text and not front.images and not back.text and not back.images:
            continue
        sides.append(front)
        sides.append(back)
    return sides


def read_cards_from_notion_markdown(md_text: str, media_dir: Path) -> List[CardSide]:
    """Parse a Notion Markdown export string into card sides.

    Notion's markdown export embeds toggles as raw ``<details>`` HTML, so we
    just need to translate inline image syntax to ``<img>`` and reuse the
    HTML walker.
    """
    if not md_text:
        return []
    return read_cards_from_notion_html(_markdown_images_to_html(md_text), media_dir)


def looks_like_notion_html(text: str) -> bool:
    """Heuristic: does this text contain at least one ``<details>`` toggle?"""
    return bool(re.search(r"<details\b", text, re.IGNORECASE))


# ---------------------------------------------------------------------------
# Font loading (shared between layout measurement and rendering)
# ---------------------------------------------------------------------------
_FONT_CANDIDATES = [
    "/System/Library/Fonts/Supplemental/Arial.ttf",
    "/System/Library/Fonts/Helvetica.ttc",
    "/Library/Fonts/Arial.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/TTF/DejaVuSans.ttf",
    "C:/Windows/Fonts/arial.ttf",
    "DejaVuSans.ttf",
    "Arial.ttf",
]

_cached_font_path: Optional[str] = None


def _find_font_path() -> Optional[str]:
    global _cached_font_path
    if _cached_font_path is not None:
        return _cached_font_path
    for candidate in _FONT_CANDIDATES:
        try:
            ImageFont.truetype(candidate, 12)
            _cached_font_path = candidate
            return candidate
        except (OSError, IOError):
            continue
    return None


def load_font(size_px: int) -> ImageFont.ImageFont:
    path = _find_font_path()
    if path is None:
        return ImageFont.load_default()
    return ImageFont.truetype(path, max(1, int(size_px)))


# ReportLab is the source of truth for PDF (and PPTX) text width; we must wrap
# lines using the same metrics as drawCentredString, otherwise PIL can measure
# a line as “fitting” while the PDF font is wider and the rendered line clips
# on both sides (a centered fragment of the sentence is all that remains).
_cached_pdf_font_name: Optional[str] = None


def _ensure_pdf_font() -> str:
    """Register the deck TTF with ReportLab (once) and return the font name."""
    global _cached_pdf_font_name
    if _cached_pdf_font_name is not None:
        return _cached_pdf_font_name
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    ttf_path = _find_font_path()
    if ttf_path and ttf_path.lower().endswith(".ttf"):
        try:
            pdfmetrics.registerFont(TTFont("DeckFont", ttf_path))
            _cached_pdf_font_name = "DeckFont"
        except Exception as exc:
            print(
                f"warning: falling back to Helvetica for layout/PDF ({exc})",
                file=sys.stderr,
            )
            _cached_pdf_font_name = "Helvetica"
    else:
        _cached_pdf_font_name = "Helvetica"
    return _cached_pdf_font_name


def _string_width_pt(text: str, size_pt: int) -> float:
    if not text:
        return 0.0
    from reportlab.pdfbase import pdfmetrics

    return float(pdfmetrics.stringWidth(text, _ensure_pdf_font(), size_pt))


# ---------------------------------------------------------------------------
# Text layout: wrap + shrink-to-fit
# ---------------------------------------------------------------------------
def _wrap_paragraph(paragraph: str, size_pt: int, max_w_pt: float) -> List[str]:
    if not paragraph:
        return [""]
    lines: List[str] = []
    current = ""
    for word in paragraph.split(" "):
        trial = word if not current else f"{current} {word}"
        if _string_width_pt(trial, size_pt) <= max_w_pt:
            current = trial
            continue

        if current:
            lines.append(current)
            current = ""

        if _string_width_pt(word, size_pt) <= max_w_pt:
            current = word
            continue

        piece = ""
        for ch in word:
            if _string_width_pt(piece + ch, size_pt) > max_w_pt and piece:
                lines.append(piece)
                piece = ch
            else:
                piece += ch
        current = piece

    if current:
        lines.append(current)
    return lines


def fit_text(text: str, width_in: float, height_in: float) -> Tuple[int, List[str]]:
    if not text:
        return MAX_FONT_PT, []

    # One PDF point = 1/72" — use the same unit for wrap checks and for height.
    _ensure_pdf_font()
    max_w_pt = width_in * 72.0
    height_pt = height_in * 72.0
    paragraphs = text.split("\n")

    def _layout_at(size_pt: int) -> Tuple[List[str], float, float]:
        lines: List[str] = []
        for para in paragraphs:
            lines.extend(_wrap_paragraph(para, size_pt, max_w_pt))
        total_h = size_pt * LINE_HEIGHT_FACTOR * len(lines)
        max_w = max((_string_width_pt(ln, size_pt) for ln in lines), default=0.0)
        return lines, total_h, max_w

    for size in range(MAX_FONT_PT, MIN_FONT_PT - 1, -1):
        lines, total_h, max_w = _layout_at(size)
        if total_h <= height_pt and max_w <= max_w_pt:
            return size, lines

    lines, _, _ = _layout_at(MIN_FONT_PT)
    return MIN_FONT_PT, lines


# ---------------------------------------------------------------------------
# Image layout: stack vertically, preserve aspect, fit to region
# ---------------------------------------------------------------------------
@dataclass
class PlacedImage:
    path: Path
    x_in: float
    y_in: float
    w_in: float
    h_in: float


def layout_images(
    image_paths: List[Path], region_w: float, region_h: float
) -> List[PlacedImage]:
    specs: List[Tuple[Path, int, int]] = []
    for p in image_paths:
        try:
            with Image.open(p) as im:
                specs.append((p, im.width, im.height))
        except Exception as exc:
            print(f"warning: could not open image {p}: {exc}", file=sys.stderr)

    if not specs:
        return []

    n = len(specs)
    gap_total = IMAGE_STACK_GAP_IN * (n - 1)
    avail_h = max(0.01, region_h - gap_total)
    slot_h = avail_h / n

    sized: List[Tuple[Path, float, float]] = []
    for p, iw, ih in specs:
        if iw <= 0 or ih <= 0:
            continue
        aspect = iw / ih
        w = slot_h * aspect
        h = slot_h
        if w > region_w:
            w = region_w
            h = region_w / aspect
        sized.append((p, w, h))

    if not sized:
        return []

    used_h = sum(h for _, _, h in sized) + IMAGE_STACK_GAP_IN * (len(sized) - 1)
    y_cursor = max(0.0, (region_h - used_h) / 2.0)

    placed: List[PlacedImage] = []
    for p, w, h in sized:
        x = (region_w - w) / 2.0
        placed.append(PlacedImage(p, x, y_cursor, w, h))
        y_cursor += h + IMAGE_STACK_GAP_IN
    return placed


# ---------------------------------------------------------------------------
# Slide layout: combine text + image regions
# ---------------------------------------------------------------------------
@dataclass
class SlideLayout:
    lines: List[str]
    font_pt: int
    text_region: Optional[Tuple[float, float, float, float]]
    images: List[PlacedImage]


def build_slide_layout(side: CardSide) -> SlideLayout:
    safe_x = SAFE_MARGIN_IN
    safe_y = SAFE_MARGIN_IN
    safe_w = SLIDE_WIDTH_IN - 2 * SAFE_MARGIN_IN
    safe_h = SLIDE_HEIGHT_IN - 2 * SAFE_MARGIN_IN

    has_text = bool(side.text)
    has_images = bool(side.images)

    img_region: Optional[Tuple[float, float, float, float]] = None
    text_region: Optional[Tuple[float, float, float, float]] = None

    if has_images and has_text:
        img_h = safe_h * IMAGE_REGION_FRAC
        txt_h = safe_h * TEXT_REGION_FRAC
        img_region = (safe_x, safe_y, safe_w, img_h)
        text_region = (safe_x, safe_y + img_h, safe_w, txt_h)
    elif has_images:
        img_region = (safe_x, safe_y, safe_w, safe_h)
    else:
        text_region = (safe_x, safe_y, safe_w, safe_h)

    placed_images: List[PlacedImage] = []
    if img_region is not None:
        ix, iy, iw, ih = img_region
        for rel in layout_images(side.images, iw, ih):
            placed_images.append(
                PlacedImage(rel.path, ix + rel.x_in, iy + rel.y_in, rel.w_in, rel.h_in)
            )

    if has_text and text_region is not None:
        _, _, tw, th = text_region
        font_pt, lines = fit_text(side.text, tw, th)
    else:
        font_pt, lines = MIN_FONT_PT, []

    return SlideLayout(
        lines=lines,
        font_pt=font_pt,
        text_region=text_region if has_text else None,
        images=placed_images,
    )


# ---------------------------------------------------------------------------
# PDF backend (reportlab)
# ---------------------------------------------------------------------------
def _render_pdf_to(sides: List[CardSide], sink) -> None:
    from reportlab.pdfgen import canvas
    from reportlab.lib.units import inch

    font_name = _ensure_pdf_font()
    page_size = (SLIDE_WIDTH_IN * inch, SLIDE_HEIGHT_IN * inch)
    c = canvas.Canvas(sink, pagesize=page_size)

    for side in sides:
        layout = build_slide_layout(side)

        c.setFillColorRGB(1, 1, 1)
        c.rect(0, 0, page_size[0], page_size[1], fill=1, stroke=0)

        for pi in layout.images:
            y_from_bottom_in = SLIDE_HEIGHT_IN - pi.y_in - pi.h_in
            try:
                c.drawImage(
                    str(pi.path),
                    pi.x_in * inch,
                    y_from_bottom_in * inch,
                    width=pi.w_in * inch,
                    height=pi.h_in * inch,
                    preserveAspectRatio=True,
                    mask="auto",
                )
            except Exception as exc:
                print(
                    f"warning: could not draw image {pi.path}: {exc}",
                    file=sys.stderr,
                )

        if layout.lines and layout.text_region is not None:
            tx, ty, tw, th = layout.text_region
            line_height_pt = layout.font_pt * LINE_HEIGHT_FACTOR
            total_h_in = line_height_pt * len(layout.lines) / 72.0
            block_top_in = ty + max(0.0, (th - total_h_in) / 2.0)

            c.setFillColorRGB(0, 0, 0)
            c.setFont(font_name, layout.font_pt)

            for i, line in enumerate(layout.lines):
                baseline_from_top_in = (
                    block_top_in + (i * line_height_pt + layout.font_pt * 0.8) / 72.0
                )
                y_from_bottom_in = SLIDE_HEIGHT_IN - baseline_from_top_in
                c.drawCentredString(
                    (tx + tw / 2.0) * inch,
                    y_from_bottom_in * inch,
                    line,
                )

        c.showPage()

    c.save()


def render_pdf(sides: List[CardSide], out_path: Path) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    _render_pdf_to(sides, str(out_path))
    return out_path


def render_pdf_bytes(sides: List[CardSide]) -> bytes:
    """Render ``sides`` to a PDF in memory and return the raw bytes."""
    buf = io.BytesIO()
    _render_pdf_to(sides, buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PPTX backend (python-pptx)
# ---------------------------------------------------------------------------
def _build_pptx(sides: List[CardSide]):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank_layout = prs.slide_layouts[6]  # "Blank"

    for side in sides:
        layout = build_slide_layout(side)
        slide = prs.slides.add_slide(blank_layout)

        bg_fill = slide.background.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        for pi in layout.images:
            try:
                # Only pass width; python-pptx derives height preserving aspect.
                slide.shapes.add_picture(
                    str(pi.path),
                    Inches(pi.x_in),
                    Inches(pi.y_in),
                    width=Inches(pi.w_in),
                )
            except Exception as exc:
                print(
                    f"warning: could not embed image {pi.path}: {exc}",
                    file=sys.stderr,
                )

        if layout.lines and layout.text_region is not None:
            tx, ty, tw, th = layout.text_region
            tb = slide.shapes.add_textbox(
                Inches(tx), Inches(ty), Inches(tw), Inches(th)
            )
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0

            for idx, line in enumerate(layout.lines):
                para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                para.alignment = PP_ALIGN.CENTER
                run = para.add_run()
                run.text = line
                run.font.size = Pt(layout.font_pt)
                run.font.color.rgb = RGBColor(0, 0, 0)

    return prs


def render_pptx(sides: List[CardSide], out_path: Path) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    _build_pptx(sides).save(str(out_path))
    return out_path


def render_pptx_bytes(sides: List[CardSide]) -> bytes:
    """Render ``sides`` to a .pptx in memory and return the raw bytes."""
    buf = io.BytesIO()
    _build_pptx(sides).save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PNG backend (Pillow)
# ---------------------------------------------------------------------------
def _render_one_png(side: CardSide) -> Image.Image:
    img_w = int(round(SLIDE_WIDTH_IN * PNG_DPI))
    img_h = int(round(SLIDE_HEIGHT_IN * PNG_DPI))
    layout = build_slide_layout(side)
    canvas_img = Image.new("RGB", (img_w, img_h), "white")

    for pi in layout.images:
        try:
            with Image.open(pi.path) as src:
                src = src.convert("RGBA")
                target_w = max(1, int(round(pi.w_in * PNG_DPI)))
                target_h = max(1, int(round(pi.h_in * PNG_DPI)))
                src.thumbnail((target_w, target_h), Image.LANCZOS)
                box_x = int(round(pi.x_in * PNG_DPI))
                box_y = int(round(pi.y_in * PNG_DPI))
                px = box_x + (target_w - src.width) // 2
                py = box_y + (target_h - src.height) // 2
                canvas_img.paste(src, (px, py), src)
        except Exception as exc:
            print(
                f"warning: could not render image {pi.path}: {exc}",
                file=sys.stderr,
            )

    if layout.lines and layout.text_region is not None:
        tx, ty, tw, th = layout.text_region
        draw = ImageDraw.Draw(canvas_img)
        font_px = max(1, int(round(layout.font_pt * PNG_DPI / 72.0)))
        font = load_font(font_px)
        line_height_px = font_px * LINE_HEIGHT_FACTOR
        total_h_px = line_height_px * len(layout.lines)
        region_top_px = ty * PNG_DPI
        region_h_px = th * PNG_DPI
        start_y_px = region_top_px + max(0.0, (region_h_px - total_h_px) / 2.0)

        region_left_px = tx * PNG_DPI
        region_w_px = tw * PNG_DPI

        for i, line in enumerate(layout.lines):
            bbox = font.getbbox(line)
            line_w = bbox[2] - bbox[0]
            x_draw = region_left_px + (region_w_px - line_w) / 2.0 - bbox[0]
            y_draw = start_y_px + i * line_height_px
            draw.text((x_draw, y_draw), line, fill="black", font=font)

    return canvas_img


def render_png(sides: List[CardSide], out_dir: Path) -> Path:
    out_dir.mkdir(parents=True, exist_ok=True)
    digits = max(3, len(str(len(sides))))
    for idx, side in enumerate(sides, start=1):
        img = _render_one_png(side)
        img.save(out_dir / f"slide_{idx:0{digits}d}.png", "PNG")
    return out_dir


def render_png_zip_bytes(sides: List[CardSide], stem: str = "slides") -> bytes:
    """Render each side to a PNG and return a ZIP archive of them as bytes."""
    buf = io.BytesIO()
    digits = max(3, len(str(len(sides))))
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for idx, side in enumerate(sides, start=1):
            img = _render_one_png(side)
            png_buf = io.BytesIO()
            img.save(png_buf, "PNG")
            zf.writestr(f"{stem}/slide_{idx:0{digits}d}.png", png_buf.getvalue())
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Anki output: pair the flat sides list back into (front, back) cards.
# ---------------------------------------------------------------------------
def sides_to_cards(sides: List[CardSide]) -> List[Tuple[CardSide, CardSide]]:
    """Group ``[front1, back1, front2, back2, …]`` into ``[(front1, back1), …]``.

    Stray trailing entries (an odd-length list) are paired with an empty
    side so nothing is silently dropped.
    """
    cards: List[Tuple[CardSide, CardSide]] = []
    for i in range(0, len(sides), 2):
        front = sides[i]
        back = sides[i + 1] if i + 1 < len(sides) else CardSide(text="", images=[])
        cards.append((front, back))
    return cards


def _side_to_anki_html(side: CardSide) -> str:
    """Render a CardSide as Anki-friendly HTML (text + ``<img>`` references).

    Images are referenced by basename so Anki resolves them against the
    deck's ``collection.media``. The text portion preserves line breaks via
    ``<br>`` tags (Anki strips raw newlines on display).
    """
    text_html = html.escape(side.text or "").replace("\n", "<br>")
    img_html = "".join(
        f'<img src="{html.escape(p.name, quote=True)}">' for p in side.images
    )
    if text_html and img_html:
        return f"{img_html}<br>{text_html}"
    return img_html or text_html


_ANKI_TXT_TAG_RE = re.compile(r"\t|\r|\n")


def _side_to_anki_txt_field(side: CardSide) -> str:
    """Render a CardSide as a single Anki TSV field with HTML preserved."""
    fragments: List[str] = []
    for p in side.images:
        fragments.append(f'<img src="{html.escape(p.name, quote=True)}">')
    if side.text:
        text_html = html.escape(side.text).replace("\n", "<br>")
        fragments.append(text_html)
    field = "<br>".join(f for f in fragments if f)
    return _ANKI_TXT_TAG_RE.sub(" ", field)


def _stable_anki_id(seed: str) -> int:
    """Derive a stable 32-bit-ish integer id from a string seed.

    Anki note/model/deck ids must be unique 64-bit integers. Re-deriving
    the same id from the same input lets a user re-export a deck and have
    Anki recognise it as the same deck on import (otherwise it'd duplicate).
    """
    import hashlib
    h = hashlib.sha256(seed.encode("utf-8")).digest()
    return int.from_bytes(h[:6], "big") | 0x100000000  # ensure > 32 bits


def render_apkg_bytes(sides: List[CardSide], deck_name: str = "Deck") -> bytes:
    """Build an Anki ``.apkg`` from ``sides`` and return the raw bytes.

    Uses the ``genanki`` library. Each pair of sides becomes one Basic card
    (front + back); referenced images are added to the package's media
    list so they ship inside the bundle.
    """
    try:
        import genanki  # lazy: only needed for this output format
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError(
            "the `genanki` package is required to build .apkg files. "
            "Install it with `pip install genanki`."
        ) from exc
    import tempfile

    deck_name = (deck_name or "Deck").strip() or "Deck"
    model_id = _stable_anki_id("anki-to-slides::model::v1")
    deck_id = _stable_anki_id(f"anki-to-slides::deck::{deck_name}")

    model = genanki.Model(
        model_id,
        "Anki to Slides Basic",
        fields=[{"name": "Front"}, {"name": "Back"}],
        templates=[{
            "name": "Card 1",
            "qfmt": "{{Front}}",
            "afmt": '{{FrontSide}}<hr id="answer">{{Back}}',
        }],
        css=(
            ".card { font-family: -apple-system, BlinkMacSystemFont, "
            "'Segoe UI', Roboto, sans-serif; font-size: 20px; "
            "text-align: center; color: #111; background: #fff; } "
            ".card img { max-width: 100%; height: auto; }"
        ),
    )

    deck = genanki.Deck(deck_id, deck_name)
    media_paths: List[str] = []
    seen_media: set = set()

    for front, back in sides_to_cards(sides):
        for side in (front, back):
            for p in side.images:
                key = p.name
                if key in seen_media:
                    continue
                try:
                    if p.is_file():
                        media_paths.append(str(p))
                        seen_media.add(key)
                except OSError:
                    continue

        front_html = _side_to_anki_html(front)
        back_html = _side_to_anki_html(back)
        if not front_html and not back_html:
            continue
        deck.add_note(genanki.Note(model=model, fields=[front_html, back_html]))

    package = genanki.Package(deck)
    package.media_files = media_paths

    with tempfile.NamedTemporaryFile(suffix=".apkg", delete=False) as tmp:
        tmp_path = tmp.name
    try:
        package.write_to_file(tmp_path)
        with open(tmp_path, "rb") as fh:
            return fh.read()
    finally:
        try:
            Path(tmp_path).unlink()
        except OSError:
            pass


def render_anki_txt_zip_bytes(sides: List[CardSide], stem: str = "deck") -> bytes:
    """Build a ``.zip`` containing an Anki Notes-in-Plain-Text export plus media.

    The zip layout mirrors what Anki itself produces when you tick "Include
    HTML and media references" — a tab-separated ``<stem>.txt`` plus the
    referenced images at the zip root — so the same archive can be fed
    back into this tool (or imported into Anki via "Import → Notes in
    Plain Text" after extracting and pointing media at the folder).
    """
    safe_stem = stem or "deck"
    buf = io.BytesIO()
    seen_media: set = set()

    txt_buf = io.StringIO()
    writer = csv.writer(
        txt_buf, delimiter="\t", quotechar='"', quoting=csv.QUOTE_MINIMAL
    )
    for front, back in sides_to_cards(sides):
        front_field = _side_to_anki_txt_field(front)
        back_field = _side_to_anki_txt_field(back)
        if not front_field and not back_field:
            continue
        writer.writerow([front_field, back_field])

    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr(f"{safe_stem}.txt", txt_buf.getvalue())
        for side in sides:
            for p in side.images:
                if p.name in seen_media:
                    continue
                try:
                    if not p.is_file():
                        continue
                    with open(p, "rb") as fh:
                        zf.writestr(p.name, fh.read())
                    seen_media.add(p.name)
                except OSError:
                    continue
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Convenience entry point (handy for the future web backend too)
# ---------------------------------------------------------------------------
def convert(
    input_path: Path,
    media_dir: Path,
    out_path: Path,
    fmt: str,
) -> Tuple[Path, int]:
    """Parse ``input_path`` and render to ``out_path`` in the requested format.

    Returns ``(final_output_path, number_of_slides)``. The returned path is:
      - the PDF file for ``pdf``
      - the .pptx file (inside its wrapping folder) for ``pptx``
      - the folder of PNGs for ``png``
    """
    fmt = fmt.lower()
    if fmt not in SUPPORTED_FORMATS:
        raise ValueError(f"unsupported format: {fmt}")

    sides = read_cards(input_path, media_dir)
    if not sides:
        raise ValueError("no cards found in input")

    if fmt == "pdf":
        final = render_pdf(sides, out_path)
    elif fmt == "pptx":
        final = render_pptx(sides, out_path)
    elif fmt == "png":
        final = render_png(sides, out_path)
    elif fmt == "apkg":
        out_path.parent.mkdir(parents=True, exist_ok=True)
        out_path.write_bytes(render_apkg_bytes(sides, deck_name=out_path.stem))
        final = out_path
    elif fmt == "anki-txt":
        out_path.parent.mkdir(parents=True, exist_ok=True)
        out_path.write_bytes(render_anki_txt_zip_bytes(sides, stem=out_path.stem))
        final = out_path
    else:  # pragma: no cover — guarded above
        raise ValueError(f"unsupported format: {fmt}")

    return final, len(sides)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------
def _pick_default_input() -> Optional[Path]:
    """If exactly one .txt lives in ./import, use it. Otherwise return None."""
    if not DEFAULT_IMPORT_DIR.is_dir():
        return None
    candidates = sorted(DEFAULT_IMPORT_DIR.glob("*.txt"))
    if len(candidates) == 1:
        return candidates[0]
    return None


def _default_output_path(stem: str, fmt: str) -> Path:
    """Compute the default output path for a given deck name + format.

    - pdf   -> export/<stem>.pdf            (single file, no wrapping folder)
    - pptx  -> export/<stem>/<stem>.pptx    (pptx wrapped in its own folder)
    - png   -> export/<stem>/               (folder of numbered PNGs)
    """
    if fmt == "pdf":
        return DEFAULT_EXPORT_DIR / f"{stem}.pdf"
    if fmt == "pptx":
        return DEFAULT_EXPORT_DIR / stem / f"{stem}.pptx"
    if fmt == "png":
        return DEFAULT_EXPORT_DIR / stem
    if fmt == "apkg":
        return DEFAULT_EXPORT_DIR / f"{stem}.apkg"
    if fmt == "anki-txt":
        return DEFAULT_EXPORT_DIR / f"{stem}.zip"
    raise ValueError(f"unsupported format: {fmt}")


def main(argv: Optional[List[str]] = None) -> int:
    ap = argparse.ArgumentParser(
        description=(
            "Convert a tab-separated Anki export into a 16:9 slide deck. "
            "One page/slide per card side, ordered front1, back1, front2, back2."
        )
    )
    ap.add_argument(
        "input",
        nargs="?",
        default=None,
        help=(
            "Path to the tab-separated Anki export (.txt). "
            "Defaults to the single .txt file in ./import when omitted."
        ),
    )
    ap.add_argument(
        "--format",
        choices=SUPPORTED_FORMATS,
        default="pdf",
        help="Output format (default: pdf).",
    )
    ap.add_argument(
        "--media",
        default=None,
        help=(
            "Directory containing media files referenced by <img src=...>. "
            f"Defaults to {DEFAULT_ANKI_MEDIA}."
        ),
    )
    ap.add_argument(
        "--out",
        default=None,
        help=(
            "Output path. Defaults to:\n"
            "  pdf  -> export/<stem>.pdf\n"
            "  pptx -> export/<stem>/<stem>.pptx\n"
            "  png  -> export/<stem>/"
        ),
    )
    args = ap.parse_args(argv)

    if args.input:
        input_path = Path(args.input)
    else:
        found = _pick_default_input()
        if found is None:
            print(
                "error: no input specified and could not auto-pick a single .txt "
                "from ./import",
                file=sys.stderr,
            )
            return 2
        input_path = found

    if not input_path.is_file():
        print(f"error: input file not found: {input_path}", file=sys.stderr)
        return 2

    media_dir = Path(args.media) if args.media else DEFAULT_ANKI_MEDIA
    if not media_dir.is_dir():
        print(
            f"warning: media directory not found: {media_dir} "
            "(images referenced by cards will be skipped)",
            file=sys.stderr,
        )

    out_path = Path(args.out) if args.out else _default_output_path(
        input_path.stem, args.format
    )

    try:
        final, n_slides = convert(input_path, media_dir, out_path, args.format)
    except ValueError as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 1

    print(f"Wrote {n_slides} slides ({args.format}) -> {final}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
